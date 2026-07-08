"""
Agno 能力接入验证(在容器内跑,agno 仅容器有):
  docker exec ezdata-backend-dev python -m module_rag._verify_agno
覆盖:Agno readers(text/csv/json/markdown)、切分策略、EsAgnoVectorDb 全 agno 接口
      (create/insert/content_hash_exists/search[vector/keyword/hybrid]/delete/drop)。
需 ES8 + DashScope + storage(MinIO)。
"""

import sys
import uuid

from agno.knowledge.document import Document
from agno.vectordb.search import SearchType

from config.env import RagConfig, TaskLogConfig
from module_rag.processing import chunk_text, read_file
from module_rag.vector_store.agno_es_vectordb import EsAgnoVectorDb
from utils.storage_utils import storage

PASS = FAIL = 0


def check(label, cond, detail=''):
    global PASS, FAIL
    ok = bool(cond)
    PASS += ok
    FAIL += not ok
    print(f'  [{"PASS" if ok else "FAIL"}] {label}' + (f' — {detail}' if detail else ''))


def main():
    print('== 1. Agno Readers(含 pptx)==')
    cases = {
        '.txt': 'ezdata 文本读取测试。\n'.encode() * 10,
        '.csv': 'name,role\nezdata,平台\ndlt,抽取\n'.encode(),
        '.json': b'{"product":"ezdata","kb":"rag"}',
        '.md': '# 标题\n\nezdata markdown 读取。'.encode(),
    }
    # 现造一个 pptx(python-pptx)
    try:
        import io as _io

        from pptx import Presentation

        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = 'ezdata pptx 抽取测试'
        buf = _io.BytesIO()
        prs.save(buf)
        cases['.pptx'] = buf.getvalue()
    except Exception as e:
        print('  (跳过 pptx 造样本:', str(e)[:50], ')')
    for ext, data in cases.items():
        key = f'upload/_agno_rd{ext}'
        storage.save(key, data)
        txt = read_file(key, filename=f'test{ext}')
        check(f'reader {ext}', len(txt) > 0, repr(txt[:30]))

    print('== 2. 切分策略(含 markdown)==')
    md_doc = '# 平台\n\nezdata 是 AI 原生数据平台。\n\n## 福利\n\n年假15天。\n\n## 技术\n\n向量库 ES8。'
    for s in ('recursive', 'fixed', 'semantic', 'markdown'):
        src = md_doc if s == 'markdown' else 'ezdata 是 AI 原生数据平台。' * 40
        cs = chunk_text(src, strategy=s, chunk_size=200, overlap=30)
        check(f'chunk {s}', len(cs) >= 1, f'{len(cs)} 段')

    print('== 3. EsAgnoVectorDb(agno VectorDb 接口)==')
    hosts = RagConfig.rag_vector_hosts or TaskLogConfig.task_es_hosts or 'http://127.0.0.1:9200'
    idx = f'rag_agno_{uuid.uuid4().hex[:8]}'
    v = EsAgnoVectorDb({'hosts': hosts}, idx, tenant_id=100, search_type=SearchType.hybrid)
    try:
        v.create()
        check('create + exists', v.exists())
        docs = [
            Document(id=uuid.uuid4().hex, name='hr', content='公司每年提供15天带薪年假', content_id='hr-1'),
            Document(
                id=uuid.uuid4().hex,
                name='tech',
                content='向量库采用 ES8 的 dense_vector 做 kNN 检索',
                content_id='tech-1',
            ),
            Document(id=uuid.uuid4().hex, name='tech', content='ETL 数据集成基于 dlt 实现', content_id='tech-2'),
        ]
        v.insert('hash-abc', docs, filters={'document_id': 'doc1'})
        check('insert 后 content_hash_exists', v.content_hash_exists('hash-abc'))
        check('id_exists', v.id_exists(docs[0].id))

        for st in (SearchType.vector, SearchType.keyword, SearchType.hybrid):
            v.search_type = st
            r = v.search('年假有多少天', limit=3)
            check(f'search {st.value} 命中', len(r) >= 1, [d.content[:12] for d in r])
            check(f'search {st.value} 返回 agno Document', all(isinstance(d, Document) for d in r))

        # 租户隔离:换租户应查不到
        v2 = EsAgnoVectorDb({'hosts': hosts}, idx, tenant_id=999, search_type=SearchType.vector)
        check('错误租户召回 0', len(v2.search('年假', limit=3)) == 0)

        check('delete_by_content_id', v.delete_by_content_id('tech-1'))
        check('delete 后计数减少', not v.id_exists(docs[1].id))
        return 0 if FAIL == 0 else 1
    finally:
        v.drop()
        for ext in cases:
            try:
                storage.delete(f'upload/_agno_rd{ext}')
            except Exception:
                pass
        print(f'\n== 结果:{PASS} passed, {FAIL} failed ==')


if __name__ == '__main__':
    sys.exit(main())
